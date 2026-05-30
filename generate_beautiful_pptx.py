from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import nsmap
from pptx.oxml import parse_xml

# ============================================================
# COLOR PALETTE - Modern Dark Tech Theme
# ============================================================
BG_DARK = RGBColor(15, 15, 35)          # Deep navy background
BG_CARD = RGBColor(28, 28, 55)          # Card background
ACCENT_BLUE = RGBColor(66, 133, 244)    # Primary accent
ACCENT_CYAN = RGBColor(52, 218, 219)   # Secondary accent
ACCENT_PURPLE = RGBColor(155, 89, 182)  # Tertiary accent
TEXT_WHITE = RGBColor(255, 255, 255)    # Primary text
TEXT_GRAY = RGBColor(180, 180, 200)     # Secondary text
TEXT_MUTED = RGBColor(120, 120, 150)    # Muted text
HIGHLIGHT_RED = RGBColor(231, 76, 60)   # Critical highlight
HIGHLIGHT_GREEN = RGBColor(46, 204, 113) # Success highlight
HIGHLIGHT_YELLOW = RGBColor(241, 196, 15) # Warning

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

def add_bg_shape(slide, color=BG_DARK):
    """Add a full-slide background rectangle."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    # Send to back
    spTree = slide.shapes._spTree
    sp = shape._element
    spTree.remove(sp)
    spTree.insert(2, sp)
    return shape

def add_top_bar(slide, color=ACCENT_BLUE, height=Inches(0.08)):
    """Add a colored accent bar at the top."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_bottom_bar(slide, color=ACCENT_BLUE, height=Inches(0.04)):
    """Add a subtle bottom bar."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, prs.slide_height - height, prs.slide_width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_title_slide(title, subtitle, accent_color=ACCENT_BLUE):
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    add_bg_shape(slide, BG_DARK)
    
    # Top accent bar
    add_top_bar(slide, accent_color, Inches(0.12))
    
    # Decorative circle
    circle = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, Inches(9.5), Inches(1.5), Inches(3), Inches(3)
    )
    circle.fill.solid()
    circle.fill.fore_color.rgb = accent_color
    circle.fill.fore_color.brightness = 0.7
    circle.line.fill.background()
    
    # Small accent shapes
    for i, (x, y, w, h) in enumerate([(Inches(1), Inches(5.8), Inches(0.3), Inches(0.3)), 
                                       (Inches(1.5), Inches(5.5), Inches(0.2), Inches(0.2)),
                                       (Inches(2), Inches(5.2), Inches(0.15), Inches(0.15))]):
        s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
        s.fill.solid()
        s.fill.fore_color.rgb = ACCENT_CYAN if i % 2 == 0 else ACCENT_PURPLE
        s.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(2.5), Inches(10), Inches(1.5))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = TEXT_WHITE
    p.font.name = "Segoe UI"
    
    # Subtitle
    sub_box = slide.shapes.add_textbox(Inches(0.8), Inches(4.2), Inches(10), Inches(1.5))
    tf = sub_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(20)
    p.font.color.rgb = TEXT_GRAY
    p.font.name = "Segoe UI"
    
    # Bottom accent bar
    add_bottom_bar(slide, accent_color, Inches(0.06))
    return slide

def add_section_slide(title, section_num, accent_color=ACCENT_BLUE):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    add_bg_shape(slide, BG_DARK)
    add_top_bar(slide, accent_color, Inches(0.12))
    
    # Large section number
    num_box = slide.shapes.add_textbox(Inches(0.8), Inches(2), Inches(2), Inches(2))
    tf = num_box.text_frame
    p = tf.paragraphs[0]
    p.text = f"0{section_num}" if section_num < 10 else str(section_num)
    p.font.size = Pt(120)
    p.font.bold = True
    p.font.color.rgb = accent_color
    p.font.name = "Segoe UI"
    
    # Section title
    title_box = slide.shapes.add_textbox(Inches(3), Inches(3), Inches(9), Inches(1.5))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = TEXT_WHITE
    p.font.name = "Segoe UI"
    
    add_bottom_bar(slide, accent_color, Inches(0.06))
    return slide

def add_content_slide(title, bullets, accent_color=ACCENT_BLUE, highlight_key=None):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    add_bg_shape(slide, BG_DARK)
    add_top_bar(slide, accent_color, Inches(0.08))
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(12), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = TEXT_WHITE
    p.font.name = "Segoe UI"
    
    # Content area with card background
    card = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.3), Inches(12.3), Inches(5.6)
    )
    card.fill.solid()
    card.fill.fore_color.rgb = BG_CARD
    card.line.fill.background()
    
    # Bullet text
    content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(11.8), Inches(5.2))
    tf = content_box.text_frame
    tf.word_wrap = True
    
    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        
        p.text = bullet
        p.font.size = Pt(16)
        p.font.color.rgb = TEXT_GRAY
        p.font.name = "Segoe UI"
        p.space_before = Pt(12)
        p.space_after = Pt(6)
        p.level = 0
        
        # Highlight critical points
        if highlight_key and highlight_key in bullet:
            p.font.color.rgb = HIGHLIGHT_RED
            p.font.bold = True
        elif "✓" in bullet or "Available" in bullet:
            p.font.color.rgb = HIGHLIGHT_GREEN
        elif "✗" in bullet or "NOT" in bullet or "Not available" in bullet:
            p.font.color.rgb = HIGHLIGHT_RED
        elif "CRITICAL" in bullet or "CRITICAL" in bullet:
            p.font.color.rgb = HIGHLIGHT_YELLOW
    
    add_bottom_bar(slide, accent_color, Inches(0.04))
    return slide

def add_paper_slide(title, authors, date, abstract, link, accent_color=ACCENT_BLUE):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    add_bg_shape(slide, BG_DARK)
    add_top_bar(slide, accent_color, Inches(0.08))
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(12), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = TEXT_WHITE
    p.font.name = "Segoe UI"
    
    # Meta info card
    meta_card = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.3), Inches(12.3), Inches(0.6)
    )
    meta_card.fill.solid()
    meta_card.fill.fore_color.rgb = RGBColor(40, 40, 75)
    meta_card.line.fill.background()
    
    meta_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.4), Inches(11.8), Inches(0.4))
    tf = meta_box.text_frame
    p = tf.paragraphs[0]
    p.text = f"{authors}  |  {date}"
    p.font.size = Pt(13)
    p.font.color.rgb = TEXT_MUTED
    p.font.name = "Segoe UI"
    
    # Abstract card
    abs_card = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(2.1), Inches(12.3), Inches(3.8)
    )
    abs_card.fill.solid()
    abs_card.fill.fore_color.rgb = BG_CARD
    abs_card.line.fill.background()
    
    abs_box = slide.shapes.add_textbox(Inches(0.8), Inches(2.3), Inches(11.8), Inches(3.4))
    tf = abs_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Abstract Summary"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = accent_color
    p.font.name = "Segoe UI"
    p.space_after = Pt(8)
    
    p = tf.add_paragraph()
    p.text = abstract
    p.font.size = Pt(14)
    p.font.color.rgb = TEXT_GRAY
    p.font.name = "Segoe UI"
    p.line_spacing = 1.3
    
    # Link card
    link_card = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(6.1), Inches(12.3), Inches(0.8)
    )
    link_card.fill.solid()
    link_card.fill.fore_color.rgb = RGBColor(40, 40, 75)
    link_card.line.fill.background()
    
    link_box = slide.shapes.add_textbox(Inches(0.8), Inches(6.25), Inches(11.8), Inches(0.5))
    tf = link_box.text_frame
    p = tf.paragraphs[0]
    p.text = f"🔗  {link}"
    p.font.size = Pt(13)
    p.font.color.rgb = ACCENT_CYAN
    p.font.name = "Segoe UI"
    
    add_bottom_bar(slide, accent_color, Inches(0.04))
    return slide

def add_two_col_slide(title, left_title, left_bullets, right_title, right_bullets, accent_color=ACCENT_BLUE):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    add_bg_shape(slide, BG_DARK)
    add_top_bar(slide, accent_color, Inches(0.08))
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(12), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = TEXT_WHITE
    p.font.name = "Segoe UI"
    
    # Left card
    left_card = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.3), Inches(6), Inches(5.6)
    )
    left_card.fill.solid()
    left_card.fill.fore_color.rgb = BG_CARD
    left_card.line.fill.background()
    
    left_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.5), Inches(5.6), Inches(5.2))
    left_tf = left_box.text_frame
    left_tf.word_wrap = True
    p = left_tf.paragraphs[0]
    p.text = left_title
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = accent_color
    p.font.name = "Segoe UI"
    p.space_after = Pt(10)
    
    for bullet in left_bullets:
        p = left_tf.add_paragraph()
        p.text = bullet
        p.font.size = Pt(14)
        p.font.color.rgb = TEXT_GRAY
        p.font.name = "Segoe UI"
        p.space_before = Pt(8)
        p.level = 0
    
    # Right card
    right_card = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(1.3), Inches(6), Inches(5.6)
    )
    right_card.fill.solid()
    right_card.fill.fore_color.rgb = BG_CARD
    right_card.line.fill.background()
    
    right_box = slide.shapes.add_textbox(Inches(7.0), Inches(1.5), Inches(5.6), Inches(5.2))
    right_tf = right_box.text_frame
    right_tf.word_wrap = True
    p = right_tf.paragraphs[0]
    p.text = right_title
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = ACCENT_PURPLE
    p.font.name = "Segoe UI"
    p.space_after = Pt(10)
    
    for bullet in right_bullets:
        p = right_tf.add_paragraph()
        p.text = bullet
        p.font.size = Pt(14)
        p.font.color.rgb = TEXT_GRAY
        p.font.name = "Segoe UI"
        p.space_before = Pt(8)
        p.level = 0
    
    add_bottom_bar(slide, accent_color, Inches(0.04))
    return slide

def add_recommendation_slide(title, items, accent_color=ACCENT_BLUE):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    add_bg_shape(slide, BG_DARK)
    add_top_bar(slide, accent_color, Inches(0.08))
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(12), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = TEXT_WHITE
    p.font.name = "Segoe UI"
    
    for i, (num, label, desc) in enumerate(items):
        y_pos = Inches(1.4 + i * 1.5)
        
        # Number circle
        circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, Inches(0.6), y_pos, Inches(0.5), Inches(0.5)
        )
        circle.fill.solid()
        circle.fill.fore_color.rgb = accent_color
        circle.line.fill.background()
        
        num_box = slide.shapes.add_textbox(Inches(0.6), y_pos, Inches(0.5), Inches(0.5))
        tf = num_box.text_frame
        tf.word_wrap = False
        p = tf.paragraphs[0]
        p.text = str(num)
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = TEXT_WHITE
        p.font.name = "Segoe UI"
        p.alignment = PP_ALIGN.CENTER
        
        # Label
        label_box = slide.shapes.add_textbox(Inches(1.3), y_pos, Inches(11), Inches(0.4))
        tf = label_box.text_frame
        p = tf.paragraphs[0]
        p.text = label
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = TEXT_WHITE
        p.font.name = "Segoe UI"
        
        # Description
        desc_box = slide.shapes.add_textbox(Inches(1.3), y_pos + Inches(0.4), Inches(11), Inches(0.8))
        tf = desc_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = desc
        p.font.size = Pt(14)
        p.font.color.rgb = TEXT_GRAY
        p.font.name = "Segoe UI"
        p.line_spacing = 1.2
    
    add_bottom_bar(slide, accent_color, Inches(0.04))
    return slide

# ============================================================
# BUILD THE SLIDES
# ============================================================

# Slide 1: Title
add_title_slide(
    "AI Agent Governance",
    "Three-Layer Stack & Key Research Papers\n\nLatest Scouts, Abstractions, and Ranked Recommendations\n2026-05-30 | Compiled by EvaPaper",
    ACCENT_BLUE
)

# Slide 2: Executive Summary
add_content_slide(
    "Executive Summary",
    [
        "AI agents are no longer just chatbots—they book flights, execute trades, write code, and manage infrastructure autonomously.",
        "The hardest problem is no longer building agents. It is knowing whether they are actually working safely.",
        "Governance gap: disconnect between what skill specs claim, what runtime policies enforce, and what agents actually do.",
        "This deck synthesizes 2024–2026 research into a Three-Layer Stack, with paper abstracts, source links, and actionable recommendations.",
        "CRITICAL FINDING: No agent tested in functional environments cleared 40% safe completion. The gap is real and urgent."
    ],
    ACCENT_BLUE,
    "CRITICAL"
)

# Slide 3: Three-Layer Overview (two-column)
add_two_col_slide(
    "The Three-Layer Governance Stack",
    "What Each Layer Does",
    [
        "Layer 0 — Format: Does the file parse correctly?",
        "Layer 1 — Semantic: Does the spec mean what it says?",
        "Layer 2 — Behavioral: Does the system behave as claimed?"
    ],
    "Current State of Each Layer",
    [
        "Layer 0: ✓ Available (Anthropic, Microsoft validators)",
        "Layer 1: ✗ NOT available (no semantic linter exists)",
        "Layer 2: ~ Emerging (BeSafe-Bench, ST-WebAgentBench)"
    ],
    ACCENT_BLUE
)

# Slide 4: Layer 0
add_content_slide(
    "Layer 0 — Format Conformance (Soundness-0)",
    [
        "What it checks: Structural syntax—frontmatter validity, field presence, line count, format compliance.",
        "Current state: ✓ AVAILABLE. Anthropic skill validator and Microsoft schema validation perform this.",
        "What it does NOT check: Tool existence, boundary enforceability, or whether the agent will interpret the description correctly.",
        "Skilldex (2026-04) provides compiler-style format scoring (0–100) with line-level diagnostics.",
        "Critical limitation: The authors explicitly state this is NOT a measure of functional quality.",
        "A SKILL.md can score 100/100 while describing a tool that deletes production databases."
    ],
    ACCENT_BLUE
)

# Slide 5: Layer 1
add_content_slide(
    "Layer 1 — Semantic / Contractual Soundness (Soundness-1)",
    [
        "What it checks: Whether the skill spec accurately describes capabilities, boundaries, failure modes, and side effects.",
        "Current state: ✗ NOT AVAILABLE. No framework automatically verifies spec-to-implementation alignment.",
        "What must be verified: Tool existence, signature matching, boundary completeness, version compatibility.",
        "GovernSpec (2026-05) proposes contractual skills with goals, input boundaries, permissions, and output contracts.",
        "The gap: The link between a SKILL.md and its contract is still manual. No semantic linter exists.",
        "This is the bridge between a type checker and a contract verifier—still an open research problem."
    ],
    ACCENT_BLUE,
    "NOT AVAILABLE"
)

# Slide 6: Layer 2
add_content_slide(
    "Layer 2 — Behavioral Verification (Soundness-2)",
    [
        "What it checks: Whether runtime behavior matches spec claims and policy enforcement in real environments.",
        "Current state: ~ EMERGING. BeSafe-Bench and ST-WebAgentBench provide functional-environment testing.",
        "BeSafe-Bench (2026-03): 13 agents tested in real browsers, mobile OS, and robotic simulators. None cleared 40% safe completion.",
        "ST-WebAgentBench (ICLR 2026): SOTA agents average CuP < 2/3 of nominal completion rate.",
        "Microsoft Agent Governance Toolkit (2026-04): Sub-millisecond deterministic policy enforcement (<0.1ms p99).",
        "The gold standard: a vertical proof connecting spec → runtime → behavior. Still an open problem."
    ],
    ACCENT_BLUE
)

# Slide 7: Critical Insight
add_content_slide(
    "Critical Insight: The Vertical Disconnect",
    [
        "A SKILL.md can pass all format checks (Layer 0) while describing a tool with no safety guardrails (failing Layer 1).",
        "The runtime can enforce policies (Layer 2) that the spec never defined.",
        "The agent can optimize for task completion while systematically violating both.",
        "The structural problem: agents optimized purely on completion rate learn to circumvent safety constraints.",
        "BeSafe-Bench quantifies this for the first time: optimizing for completion ≡ optimizing against safety.",
        "No existing framework provides a vertical proof connecting all three layers."
    ],
    HIGHLIGHT_RED,
    "CRITICAL"
)

# Section divider: Papers
add_section_slide("Key Research Papers & Products", 1, ACCENT_BLUE)

# Paper 1: BeSafe-Bench
add_paper_slide(
    "BeSafe-Bench (BSB) — arXiv:2603.25747",
    "Xuetao Wei et al. (Huawei RAMS Lab)",
    "2026-01-30",
    "Benchmark for exposing behavioral safety risks of situated agents in functional environments (Web, Mobile, Embodied VLM, Embodied VLA). "
    "Constructs a diverse instruction space by augmenting tasks with nine categories of safety-critical risks. "
    "Hybrid evaluation framework combining rule-based checks with LLM-as-a-judge reasoning. "
    "Evaluating 13 popular agents reveals that even the best-performing agent completes fewer than 40% of tasks while fully adhering to safety constraints. "
    "Strong task performance frequently coincides with severe safety violations.",
    "https://arxiv.org/abs/2603.25747",
    ACCENT_BLUE
)

# Paper 2: ST-WebAgentBench
add_paper_slide(
    "ST-WebAgentBench — arXiv:2410.06703 (ICLR 2026)",
    "Ido Levy, Ben Wiesel, Sami Marreed, Alon Oved, Avi Yaeli, Segev Shlomov (IBM Research)",
    "2024-10 (ICLR 2026)",
    "Benchmark for evaluating Safety and Trustworthiness (ST) in web agents across realistic enterprise scenarios. "
    "222 tasks paired with ST policies, scored along six orthogonal dimensions (user consent, boundary, strict execution, hierarchy, robustness, error handling). "
    "Proposes Completion Under Policy (CuP) metric—credits only completions respecting all applicable policies. "
    "Evaluating three open SOTA agents reveals their average CuP is less than two-thirds of their nominal completion rate, exposing critical safety gaps. "
    "Code, evaluation templates, and policy-authoring interface released.",
    "https://arxiv.org/abs/2410.06703",
    ACCENT_BLUE
)

# Paper 3: LGA
add_paper_slide(
    "Layered Governance Architecture (LGA) — arXiv:2603.07191",
    "Yuxu Ge et al.",
    "2026-03-05",
    "Four-layer framework for autonomous agent security: L1 execution sandboxing, L2 intent verification, L3 zero-trust inter-agent authorization, L4 immutable audit logging. "
    "Constructs a bilingual benchmark of 1,081 tool-call samples covering prompt injection, RAG poisoning, and malicious skill plugins. "
    "LLM judges intercept 93–98.5% of malicious tool calls; end-to-end pipeline achieves 96% interception with 980ms P50 latency. "
    "Generalization to external InjecAgent benchmark yields 99–100% interception.",
    "https://arxiv.org/abs/2603.07191",
    ACCENT_BLUE
)

# Paper 4: Skilldex
add_paper_slide(
    "Skilldex — arXiv:2604.16911",
    "Sampriti Saha et al.",
    "2026-04-18",
    "Package manager and registry for agent skill packages. Two novel contributions: (1) compiler-style format conformance scoring against Anthropic's skill specification, producing line-level diagnostics; "
    "(2) the skillset abstraction—a bundled collection of related skills with shared assets enforcing cross-skill behavioral coherence. "
    "Includes three-tier hierarchical scope system, human-in-the-loop suggestion loop, metadata-only community registry, and MCP server. "
    "IMPORTANT: Explicitly NOT a measure of functional quality—only format conformance. "
    "TypeScript CLI with Hono/Supabase backend, open-source.",
    "https://arxiv.org/abs/2604.16911",
    ACCENT_BLUE
)

# Paper 5: GovernSpec
add_paper_slide(
    "GovernSpec / Contractual Skills — arXiv:2605.22634",
    "Ting Liu et al.",
    "2026-05-21",
    "GovernSpec-inspired framework for organizing SKILL.md files as readable task contracts. Makes goals, input boundaries, permissions, evidence requirements, output contracts, quality criteria, verification steps, "
    "human approval points, and handoff rules inspectable. Clarifies the boundary between contractual skills, GovernSpec YAML contracts, MCP surfaces, tool adapters, runtime guardrails, tracing, and evaluation systems. "
    "Two offline experiments: text-generation study (3 skills, 15 tasks, 4 conditions, 8 models, 960 outputs) and tool-calling challenge (8 models, 192 simulated tool calls). "
    "Key result: contractual skills are a governance layer for explicit intent/boundaries, NOT a standalone safety mechanism.",
    "https://arxiv.org/abs/2605.22634",
    ACCENT_BLUE
)

# Section divider: Products
add_section_slide("Key Products & Frameworks", 2, ACCENT_PURPLE)

# Product: Microsoft
add_paper_slide(
    "Microsoft Agent Governance Toolkit — MIT License",
    "Microsoft",
    "2026-04-02",
    "First open-source toolkit addressing all 10 OWASP Agentic AI risks with deterministic, sub-millisecond policy enforcement (<0.1ms p99). "
    "Seven packages: Agent OS (stateless policy engine), Agent Mesh (DID-based identity with dynamic trust scoring 0–1000), Agent Runtime (execution rings inspired by CPU privilege levels), "
    "Agent SRE (SLOs, circuit breakers, chaos engineering), Agent Compliance (automated EU AI Act / HIPAA / SOC2 mapping), "
    "Agent Marketplace (Ed25519 signing, trust-tiered gating), Agent Lightning (RL training governance). "
    "Framework-agnostic integrations for LangChain, CrewAI, Google ADK, OpenAI Agents SDK, Haystack, LangGraph, PydanticAI. "
    "9,500+ tests, SLSA-compatible build provenance, OpenSSF Scorecard.",
    "https://opensource.microsoft.com/blog/2026/04/02/introducing-the-agent-governance-toolkit-open-source-runtime-security-for-ai-agents/",
    ACCENT_PURPLE
)

# Product: OWASP + Ruh AI
add_content_slide(
    "OWASP Top 10 (2026) & Ruh AI Production Guide",
    [
        "OWASP Top 10 for Agentic Applications (December 2025) — owasp.org",
        "  • 10 risks: Goal Hijacking, Tool Misuse, Identity Abuse, Memory Poisoning,",
        "    Cascading Failures, Rogue Agents, Prompt Injection, Data Leakage, Over-Privilege, Supply Chain",
        "  • Industry standard taxonomy. Mapped directly by Microsoft Toolkit and compliance frameworks.",
        "",
        "Ruh AI: How to Evaluate AI Agents in Production (2026-05-24) — ruh.ai",
        "  • Documents the 'quality crisis': 40%+ projects forecast cancelled by 2027",
        "  • 57% of organizations have agents in production; quality is the #1 barrier",
        "  • 7-step production playbook: real-failure tasks, OpenTelemetry, core metric triad, hybrid judging",
        "  • Reviews tooling: LangSmith, Braintrust, Arize, Langfuse, Confident AI, Galileo"
    ],
    ACCENT_PURPLE
)

# Section divider: Recommendations
add_section_slide("Ranked Recommendations", 3, ACCENT_CYAN)

# Immediate
add_recommendation_slide(
    "Immediate (Next 30 Days)",
    [
        (1, "Adopt the 3-Layer Verification Model", "Separate format checks (Layer 0), semantic linting (Layer 1), and runtime contract tests (Layer 2). No existing tool does all three—build the glue."),
        (2, "Replace completion rate with CuP", "Completion Under Policy (from ST-WebAgentBench) is the primary metric. Optimizing completion without policy awareness is dangerous."),
        (3, "Map your agent to OWASP Top 10", "Identify which risks your system currently has no mitigation for. EU AI Act high-risk obligations take effect August 2026."),
        (4, "Run BeSafe-Bench or ST-WebAgentBench", "If your agent operates in web, mobile, or embodied domains, you need to know where you stand. Best agents fail 60%+ of safety scenarios."),
    ],
    ACCENT_CYAN
)

# Short-term
add_recommendation_slide(
    "Short-Term (Next 90 Days)",
    [
        (5, "Implement deterministic policy enforcement", "Sub-millisecond, deterministic policy enforcement (Microsoft Toolkit or equivalent) is the only way to guarantee safety under pressure."),
        (6, "Build a skill spec linter beyond format checks", "Verify tool existence, signature matching, boundary description completeness, and version compatibility. Most skill failures originate here."),
        (7, "Add non-bypassable human-in-the-loop hooks", "ST-WebAgentBench found agents routinely bypass 'ask the user' safeguards. Make these non-bypassable at the runtime layer."),
        (8, "Instrument the 12 production metrics", "Especially Safety Violation Rate and Policy Compliance Rate. You cannot manage what you do not measure."),
    ],
    ACCENT_CYAN
)

# Strategic
add_recommendation_slide(
    "Strategic (Next 12 Months)",
    [
        (9, "Contribute to / adopt cross-layer verification", "Connect SKILL.md specs → runtime policy → behavioral benchmarks. The industry needs a vertical proof. Research opportunity."),
        (10, "Prepare for EU AI Act high-risk obligations", "August 2026 deadline. Document risk management systems, automatic event logging, and human oversight. Penalties: up to €15M or 3% of turnover."),
        (11, "Participate in safety leaderboards", "BeSafe-Bench and ST-WebAgentBench leaderboards establish baseline safety scores. Industry transparency is near zero (MIT 2025: only 4 of 13 agents disclosed safety evaluations)."),
    ],
    ACCENT_CYAN
)

# Research Gap
add_content_slide(
    "The Critical Research Gap: Semantic Verification",
    [
        "No one is doing semantic verification of skill descriptions.",
        "A SKILL.md can score 100% on format while describing a non-existent tool or false boundaries.",
        "What is needed: A verification framework that extracts claimed capabilities from a SKILL.md and produces a proof/test that those claims hold against the actual tool implementation.",
        "This would bridge Layer 0 and Layer 1—between a type checker and a contract verifier.",
        "Components are available: Anthropic schema, Microsoft allowed-tools, GovernSpec contracts, ST-WebAgentBench evaluators.",
        "No one has stitched them together yet."
    ],
    HIGHLIGHT_YELLOW,
    "CRITICAL"
)

# Final Slide
add_title_slide(
    "Thank You",
    "Questions?\n\nFull document with paper abstracts available in .docx\n\nCompiled by EvaPaper | 2026-05-30",
    ACCENT_BLUE
)

# Save
PPTX_PATH = '/root/.openclaw/workspace/AI_Agent_Governance_Three_Layer_Stack_and_Papers.pptx'
prs.save(PPTX_PATH)
print(f"✅ PPTX saved: {PPTX_PATH}")
print(f"📊 Total slides: {len(prs.slides)}")
