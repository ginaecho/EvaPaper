from __future__ import annotations

from dataclasses import dataclass
from pathlib import Path
import re

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parent

INPUTS = [
    ROOT / "AI_Agent_Governance_Three_Layer_Stack_and_Papers.md",
    ROOT / "Agent_Governance_Three_Questions_Synthesis.md",
]

SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

BG = RGBColor(12, 19, 33)
CARD = RGBColor(24, 36, 58)
CARD_ALT = RGBColor(30, 47, 73)
TEXT = RGBColor(245, 247, 250)
TEXT_SOFT = RGBColor(195, 205, 220)
ACCENTS = [
    RGBColor(63, 142, 252),
    RGBColor(62, 207, 185),
    RGBColor(166, 110, 255),
    RGBColor(255, 179, 71),
]


@dataclass
class Block:
    kind: str
    level: int = 0
    text: str = ""
    items: list[str] | None = None


def clean_inline(text: str) -> str:
    text = text.strip()
    text = re.sub(r"`([^`]+)`", r"\1", text)
    text = re.sub(r"\*\*([^*]+)\*\*", r"\1", text)
    text = re.sub(r"\*([^*]+)\*", r"\1", text)
    text = re.sub(r"\[([^\]]+)\]\(([^)]+)\)", r"\1", text)
    text = re.sub(r"\s+", " ", text)
    return text.strip()


def shorten(text: str, limit: int = 180) -> str:
    text = clean_inline(text)
    if len(text) <= limit:
        return text
    return text[: limit - 1].rstrip() + "…"


def parse_markdown(path: Path) -> list[Block]:
    lines = path.read_text(encoding="utf-8").splitlines()
    blocks: list[Block] = []
    i = 0
    while i < len(lines):
        stripped = lines[i].strip()
        if not stripped:
            i += 1
            continue
        if stripped.startswith("#"):
            level = len(stripped) - len(stripped.lstrip("#"))
            blocks.append(Block("heading", level=level, text=clean_inline(stripped[level:].strip())))
            i += 1
            continue
        if re.match(r"^([-*]|\d+\.)\s+", stripped):
            items = []
            while i < len(lines):
                candidate = lines[i].strip()
                if not re.match(r"^([-*]|\d+\.)\s+", candidate):
                    break
                items.append(clean_inline(re.sub(r"^([-*]|\d+\.)\s+", "", candidate)))
                i += 1
            blocks.append(Block("list", items=items))
            continue
        if stripped.startswith(">"):
            quote_lines = []
            while i < len(lines) and lines[i].strip().startswith(">"):
                quote_lines.append(clean_inline(lines[i].strip().lstrip(">").strip()))
                i += 1
            blocks.append(Block("paragraph", text=" ".join(quote_lines)))
            continue
        if stripped.startswith("|"):
            table_lines = []
            while i < len(lines) and lines[i].strip().startswith("|"):
                row = [clean_inline(cell) for cell in lines[i].strip().strip("|").split("|")]
                if not all(set(cell) <= {"-", ":"} for cell in row):
                    table_lines.append(" | ".join(row))
                i += 1
            if table_lines:
                blocks.append(Block("list", items=table_lines[:5]))
            continue
        para = [stripped]
        i += 1
        while i < len(lines):
            candidate = lines[i].strip()
            if (
                not candidate
                or candidate.startswith("#")
                or candidate.startswith(">")
                or candidate.startswith("|")
                or re.match(r"^([-*]|\d+\.)\s+", candidate)
            ):
                break
            para.append(candidate)
            i += 1
        blocks.append(Block("paragraph", text=clean_inline(" ".join(para))))
    return blocks


def set_bg(slide, color=BG):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    sp_tree = slide.shapes._spTree
    sp = shape._element
    sp_tree.remove(sp)
    sp_tree.insert(2, sp)


def add_top_accent(slide, color):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_WIDTH, Inches(0.12))
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()


def add_footer(slide, note: str, color):
    footer = slide.shapes.add_textbox(Inches(0.5), Inches(7.0), Inches(12.2), Inches(0.25))
    tf = footer.text_frame
    p = tf.paragraphs[0]
    p.text = note
    p.font.size = Pt(9)
    p.font.color.rgb = color
    p.alignment = PP_ALIGN.RIGHT


def add_textbox(slide, left, top, width, height, text, size, color=TEXT, bold=False, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = size
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = "Aptos"
    p.alignment = align
    return box


def add_card(slide, left, top, width, height, fill=CARD):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = fill
    return shape


def first_heading(blocks: list[Block], level: int) -> str | None:
    for block in blocks:
        if block.kind == "heading" and block.level == level:
            return block.text
    return None


def collect_summary_points(blocks: list[Block], max_points: int = 5) -> list[str]:
    points: list[str] = []
    for block in blocks:
        if block.kind == "list" and block.items:
            for item in block.items:
                if item:
                    points.append(shorten(item, 150))
                    if len(points) >= max_points:
                        return points
        elif block.kind == "paragraph" and block.text:
            points.append(shorten(block.text, 150))
            if len(points) >= max_points:
                return points
    return points


def build_sections(blocks: list[Block]):
    sections = []
    current = None
    subsection = None
    for block in blocks:
        if block.kind == "heading" and block.level == 2:
            current = {"title": block.text, "blocks": [], "subsections": []}
            sections.append(current)
            subsection = None
        elif block.kind == "heading" and block.level >= 3 and current is not None:
            subsection = {"title": block.text, "blocks": []}
            current["subsections"].append(subsection)
        elif current is not None:
            current["blocks"].append(block)
            if subsection is not None:
                subsection["blocks"].append(block)
    return sections


def extract_named_findings(subsection_blocks: list[Block], max_items: int = 6) -> list[str]:
    findings: list[str] = []
    for block in subsection_blocks:
        if block.kind == "list" and block.items:
            for item in block.items:
                cleaned = clean_inline(item)
                if "arXiv:" in cleaned:
                    findings.append(shorten(cleaned, 120))
                    if len(findings) >= max_items:
                        return findings
    return findings


def add_title_slide(prs: Presentation, title: str, subtitle: str, accent):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_top_accent(slide, accent)
    orb = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(9.7), Inches(1.0), Inches(2.3), Inches(2.3))
    orb.fill.solid()
    orb.fill.fore_color.rgb = accent
    orb.line.fill.background()
    add_textbox(slide, Inches(0.8), Inches(2.0), Inches(8.6), Inches(1.2), title, Pt(28), bold=True)
    add_textbox(slide, Inches(0.85), Inches(3.1), Inches(8.0), Inches(1.0), subtitle, Pt(15), color=TEXT_SOFT)
    add_textbox(slide, Inches(0.85), Inches(5.7), Inches(6.8), Inches(0.5), "Generated locally from the current markdown sources", Pt(11), color=TEXT_SOFT)
    add_footer(slide, "EvaPaper", accent)


def add_overview_slide(prs: Presentation, title: str, bullets: list[str], accent):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_top_accent(slide, accent)
    add_textbox(slide, Inches(0.6), Inches(0.45), Inches(8.5), Inches(0.6), title, Pt(24), bold=True)
    add_card(slide, Inches(0.55), Inches(1.3), Inches(12.2), Inches(5.7), CARD)
    box = slide.shapes.add_textbox(Inches(0.9), Inches(1.65), Inches(11.4), Inches(5.0))
    tf = box.text_frame
    tf.word_wrap = True
    for i, bullet in enumerate(bullets[:6]):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = bullet
        p.font.size = Pt(18)
        p.font.color.rgb = TEXT_SOFT
        p.font.name = "Aptos"
        p.level = 0
        p.space_after = Pt(10)
    add_footer(slide, "Overview", accent)


def add_three_card_slide(prs: Presentation, title: str, cards: list[tuple[str, list[str]]], accent):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_top_accent(slide, accent)
    add_textbox(slide, Inches(0.6), Inches(0.45), Inches(10), Inches(0.6), title, Pt(24), bold=True)
    widths = [Inches(3.9), Inches(3.9), Inches(3.9)]
    lefts = [Inches(0.55), Inches(4.7), Inches(8.85)]
    fills = [CARD, CARD_ALT, CARD]
    for idx, ((card_title, bullets), left, width, fill) in enumerate(zip(cards[:3], lefts, widths, fills)):
        add_card(slide, left, Inches(1.5), width, Inches(5.1), fill)
        add_textbox(slide, left + Inches(0.18), Inches(1.8), width - Inches(0.35), Inches(0.6), card_title, Pt(18), color=ACCENTS[idx % len(ACCENTS)], bold=True)
        box = slide.shapes.add_textbox(left + Inches(0.2), Inches(2.45), width - Inches(0.4), Inches(3.8))
        tf = box.text_frame
        tf.word_wrap = True
        for j, bullet in enumerate(bullets[:4]):
            p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
            p.text = shorten(bullet, 95)
            p.font.size = Pt(13)
            p.font.color.rgb = TEXT_SOFT
            p.font.name = "Aptos"
            p.space_after = Pt(8)
    add_footer(slide, "Illustrated summary", accent)


def add_section_divider(prs: Presentation, title: str, accent):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_top_accent(slide, accent)
    add_textbox(slide, Inches(0.8), Inches(2.2), Inches(10.5), Inches(1.0), title, Pt(30), bold=True)
    add_textbox(slide, Inches(0.82), Inches(3.35), Inches(8.5), Inches(0.6), "Curated from the current markdown report", Pt(13), color=TEXT_SOFT)
    marker = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.HEXAGON, Inches(10.4), Inches(1.8), Inches(1.8), Inches(1.8))
    marker.fill.solid()
    marker.fill.fore_color.rgb = accent
    marker.line.fill.background()
    add_footer(slide, "Section", accent)


def add_content_slide(prs: Presentation, title: str, bullets: list[str], accent):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_top_accent(slide, accent)
    add_textbox(slide, Inches(0.6), Inches(0.42), Inches(11.6), Inches(0.6), title, Pt(22), bold=True)
    add_card(slide, Inches(0.55), Inches(1.25), Inches(12.2), Inches(5.8), CARD)
    box = slide.shapes.add_textbox(Inches(0.85), Inches(1.6), Inches(11.4), Inches(5.0))
    tf = box.text_frame
    tf.word_wrap = True
    for i, bullet in enumerate(bullets[:6]):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = shorten(bullet, 180)
        p.font.size = Pt(16)
        p.font.color.rgb = TEXT_SOFT
        p.font.name = "Aptos"
        p.space_after = Pt(10)
    add_footer(slide, "EvaPaper", accent)


def make_deck(source: Path) -> Path:
    blocks = parse_markdown(source)
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    title = first_heading(blocks, 1) or source.stem
    subtitle = first_heading(blocks, 2) or "Current synthesized findings"
    summary_points = collect_summary_points(blocks, max_points=6)
    sections = build_sections(blocks)

    accent = ACCENTS[0]
    add_title_slide(prs, title, subtitle, accent)
    if summary_points:
        add_overview_slide(prs, "Key Takeaways", summary_points, ACCENTS[1])

    if "Three-Layer" in title:
        layer_cards = [
            ("Layer 0", ["Spec validation", "Contract boundaries", "Skill/package conformance", "Human approval points"]),
            ("Layer 1", ["Runtime policy", "Trust boundaries", "Inter-agent authorization", "Auditability"]),
            ("Layer 2", ["Behavioral safety", "Benchmarks", "Policy adherence", "Real-world robustness"]),
        ]
        add_three_card_slide(prs, "Governance Stack", layer_cards, ACCENTS[2])
        for section in sections:
            if section["title"] == "Executive Summary":
                for subsection in section["subsections"]:
                    if "scout addendum" in subsection["title"].lower():
                        findings = extract_named_findings(subsection["blocks"])
                        if findings:
                            add_overview_slide(prs, "June 7 Scout Addendum", findings, ACCENTS[3])
                        break
    elif "Three Critical Questions" in title:
        question_cards = [
            ("Q1", ["Do markdown skills affect behavior?", "Answer: yes, but probabilistically", "Strongest evidence: AGENTS.md and formatting studies"]),
            ("Q2", ["Can we check before runtime?", "Answer: partially yes", "Breakthrough: ZipperGen for multi-agent coordination"]),
            ("Q3", ["Can runtime checks be deterministic?", "Answer: yes, for structure and policy", "Best current formal side: ZipperGen, ACP, ABC"]),
        ]
        add_three_card_slide(prs, "Three Questions at a Glance", question_cards, ACCENTS[3])

    for idx, section in enumerate(sections[:8]):
        accent = ACCENTS[idx % len(ACCENTS)]
        add_section_divider(prs, section["title"], accent)
        subsections = section["subsections"][:4] if section["subsections"] else []
        if subsections:
            for subsection in subsections:
                bullets = collect_summary_points(subsection["blocks"], max_points=5)
                if bullets:
                    add_content_slide(prs, subsection["title"], bullets, accent)
        else:
            bullets = collect_summary_points(section["blocks"], max_points=5)
            if bullets:
                add_content_slide(prs, section["title"], bullets, accent)

    output = source.with_suffix(".pptx")
    prs.save(output)
    return output


def main() -> None:
    generated = [make_deck(path) for path in INPUTS]
    for path in generated:
        print(f"PPTX saved: {path}")


if __name__ == "__main__":
    main()
