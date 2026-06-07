#!/usr/bin/env python3
"""Convert markdown to pptx using python-pptx."""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor

def md_to_pptx(md_path, pptx_path):
    with open(md_path, 'r') as f:
        lines = f.readlines()
    
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # Title slide
    blank_layout = prs.slide_layouts[6]  # blank
    
    def add_title_slide(title, subtitle=""):
        slide = prs.slides.add_slide(blank_layout)
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.333), Inches(1.5))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(44)
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER
        if subtitle:
            p2 = tf.add_paragraph()
            p2.text = subtitle
            p2.font.size = Pt(20)
            p2.alignment = PP_ALIGN.CENTER
        return slide
    
    def add_content_slide(title, bullets):
        slide = prs.slides.add_slide(blank_layout)
        # Title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(1))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(32)
        p.font.bold = True
        
        # Content
        content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(12.333), Inches(5.5))
        tf = content_box.text_frame
        tf.word_wrap = True
        for i, bullet in enumerate(bullets):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = bullet
            p.font.size = Pt(16)
            p.space_after = Pt(8)
        return slide
    
    # Title
    add_title_slide("AI Agent Governance:\nThree-Layer Stack & Research", "Updated June 4, 2026 | 24 papers + 22 products")
    
    # Overview slide
    add_content_slide("The Three-Layer Stack", [
        "Layer 0 — Spec-Level: Govern how agents/skills are defined (SKILL.md, prompts, schemas)",
        "Layer 1 — Runtime-Level: Govern how agents execute (sandboxing, tool approval, observability)",
        "Layer 2 — Behavioral-Level: Govern what agents do (benchmarks, evaluation, red-teaming)",
        "Layer 3 — Meta-Governance: Govern the governance itself (audit, human oversight, compliance)"
    ])
    
    # Key stats slide
    add_content_slide("Why This Matters (Key Findings)", [
        "34% of community skill packages are malformed (Skilldex)",
        "89% of attacks succeed against baseline guardrails (LGA)",
        "60%+ of tasks have safety violations (BeSafe-Bench)",
        "40% of agent projects face cancellation (Agent Evaluation Guide)",
        "0% concordance across safety benchmarks (Kendall's W = 0.10)",
        "100% of individually-valid requests in attack chain approved by stateless engines",
        "1,200+ malicious skills infiltrated OpenClaw marketplace (ClawHavoc)",
        "Governance is the difference between agents that work and agents that are safe to use"
    ])
    
    # New papers overview
    add_content_slide("New Papers (June 4, 2026 Scout)", [
        "1. Taxonomy & Consistency Analysis (2605.16282) — 40 benchmarks, zero concordance",
        "2. Governance by Construction / CUGA (2605.20874) — IBM, policy-as-code, 5 checkpoints",
        "3. Agent Behavioral Contracts / ABC (2602.22302) — Design-by-Contract for agents, Drift Bounds Theorem",
        "4. AgentVerify (202604.1029) — LTL model checking, 86.67% verification accuracy",
        "5. SkillFortify (2603.00195) — Formal skill supply chain verification, F1=96.95%",
        "6. AgentAssay (2603.02601) — Regression testing for non-deterministic agents, 78-100% cost reduction",
        "7. Comprehensive Survey on Agent Skills (2605.07358) — Full skill lifecycle survey",
        "8. Machine Identity Governance Taxonomy (2604.06148) — CSA, 37 risk sub-categories",
        "9. Admission Control / ACP (2603.18829) — Temporal admission control, 1.7M req/s, TLA+ verified",
        "10. Unified Review: Memory, Skills, Protocols (2604.08224) — Externalization theory"
    ])
    
    # ABC detail
    add_content_slide("ABC: Agent Behavioral Contracts", [
        "Contract C = (P, I, G, R): Preconditions, Invariants, Governance, Recovery",
        "(p, delta, k)-satisfaction: probabilistic compliance for LLM non-determinism",
        "Drift Bounds Theorem: contracts with recovery rate gamma > alpha bound drift to D* = alpha/gamma",
        "AgentAssert runtime enforcement: <10 ms overhead per action",
        "1,980 sessions across 7 models: 5.2-6.8 soft violations detected per session",
        "88-100% hard constraint compliance; behavioral drift bounded to D* < 0.27",
        "Open source: part of Qualixar suite"
    ])
    
    # ACP detail
    add_content_slide("ACP: Admission Control for Agent Actions", [
        "Temporal (history-aware) admission control — not per-request evaluation",
        "500-request experiment: stateless engines approve 100% of individually-valid requests",
        "ACP limits autonomous execution to 2 out of 500 (0.4%)",
        "Decision latency: 739-832 ns (p50); throughput: 1,720,000 req/s",
        "TLA+ verified: 11 invariants + 4 temporal properties, 0 violations",
        "Verified across 4,294,930,695 distinct states",
        "Paper 1 of 6-paper Agent Governance Series by Marcelo Fernandez",
        "Go reference implementation: 23 packages, 138 conformance tests"
    ])
    
    # SkillFortify detail
    add_content_slide("SkillFortify: Formal Skill Supply Chain Security", [
        "First formal (not heuristic) skill supply chain verification framework",
        "DY-Skill attacker model: Dolev-Yao adapted to 5-phase skill lifecycle",
        "Sound static analysis via abstract interpretation; capability confinement proof",
        "Agent Dependency Graph with SAT-based resolution (<100 ms for 1,000 nodes)",
        "F1=96.95%, 0% false positive rate on 540-skill benchmark",
        "Replaces 'no findings does not mean no risk' with mathematical guarantees",
        "pip install skillfortify | MIT license | https://github.com/varun369/skillfortify"
    ])
    
    # Benchmark consistency
    add_content_slide("Critical Finding: Benchmarks Are Noise", [
        "Taxonomy and Consistency Analysis (2605.16282) catalogs 40 behavioral safety benchmarks",
        "Six-axis taxonomy applied across all benchmarks",
        "Coverage matrix reveals broad coverage but limited methodological convergence",
        "Kendall's W = 0.10, p = 0.94: NO ranking concordance across evaluation dimensions",
        "Conclusion: benchmark choice can yield contradictory safety conclusions",
        "Coverage counts often overstate evaluation depth",
        "Robustness remains effectively unbenchmarked",
        "Implication: NEVER trust a single benchmark. Use multiple + compare methodological axes."
    ])
    
    # New products
    add_content_slide("New Products & Frameworks", [
        "AgentAssert — Runtime enforcement library implementing ABC contracts (<10 ms overhead)",
        "SkillFortify — pip install skillfortify | Formal skill supply chain verification",
        "AgentAssay — pip install agentassay | Token-efficient regression testing for agents",
        "Agent Control Protocol (ACP) — Go ref implementation, 1.7M req/s, TLA+ verified",
        "SkillsVote — Skills engine for AI agents (ecosystem infrastructure)"
    ])
    
    # Implementation roadmap
    add_content_slide("Updated Implementation Roadmap", [
        "Immediate (30 days): Adopt ABC contract structure; deploy SkillFortify in CI/CD",
        "Short-term (90 days): Deploy AgentAssert runtime enforcement; implement ACP admission control",
        "Short-term (90 days): Integrate AgentAssay for regression testing; deploy LGA Layer 1 sandboxing",
        "Strategic (12 months): Full LGA stack; Microsoft-style governance toolkit",
        "Strategic (12 months): Skill-level attestation per MIGT; benchmark consistency dashboard",
        "Strategic (12 months): Adopt full Agent Governance Series (6 papers) as formal foundation"
    ])
    
    # Conclusion
    add_content_slide("The Shift: From Heuristic to Formal Guarantees", [
        "The question is no longer 'should we govern agents?' but 'which formal framework first?'",
        "ABC provides the specification and runtime contract foundation",
        "AgentVerify proves formal verification of observable control flow is tractable",
        "SkillFortify replaces heuristic scanning with mathematical guarantees",
        "AgentAssay makes governance continuous through regression testing",
        "ACP provides production-ready admission control with 1.7M req/s throughput",
        "The Agent Governance Series (6 papers) is the most concentrated formal governance program",
        "Benchmark inconsistency (Kendall's W=0.10) means we need meta-governance, not just more benchmarks"
    ])
    
    # Contact
    add_title_slide("EvaPaper Agent Governance Scout", "Report compiled June 4, 2026\nRepo: github.com/ginaecho/EvaPaper\nCommit: 8fb6dce")
    
    prs.save(pptx_path)
    print(f"Saved: {pptx_path}")

if __name__ == '__main__':
    md_to_pptx('/root/.openclaw/workspace/AI_Agent_Governance_Three_Layer_Stack_and_Papers.md',
                '/root/.openclaw/workspace/AI_Agent_Governance_Three_Layer_Stack_and_Papers.pptx')
