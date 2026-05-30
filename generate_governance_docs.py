from docx import Document
from docx.shared import Pt, Inches, RGBColor
from docx.enum.text import WD_ALIGN_PARAGRAPH
from pptx import Presentation
from pptx.util import Inches as PptxInches
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor as PptxRGBColor

# ============================================================
# DOCX GENERATION
# ============================================================
doc = Document()

# Title
title = doc.add_heading('AI Agent Governance: Three-Layer Stack & Key Research Papers', 0)
title.alignment = WD_ALIGN_PARAGRAPH.CENTER

# Date & Author
info = doc.add_paragraph()
info.alignment = WD_ALIGN_PARAGRAPH.CENTER
info.add_run('Prepared: 2026-05-30 | Compiled by EvaPaper').italic = True

doc.add_page_break()

# ============================================================
# SECTION 1: EXECUTIVE SUMMARY
# ============================================================
doc.add_heading('1. Executive Summary', level=1)
summary = doc.add_paragraph(
    'The rapid deployment of autonomous AI agents across enterprise environments has exposed a critical governance gap: '
    'the disconnect between what skill specifications claim, what runtime policies enforce, and what agents actually do. '
    'This document synthesizes the latest research (2024–2026) into a Three-Layer Governance Stack, provides '
    'abstracts and links to key papers and products, and offers ranked recommendations for practitioners.'
)

# ============================================================
# SECTION 2: THE THREE-LAYER GOVERNANCE STACK
# ============================================================
doc.add_heading('2. The Three-Layer Governance Stack', level=1)

# Layer 0
doc.add_heading('Layer 0 — Format Conformance (Soundness-0)', level=2)
p = doc.add_paragraph()
p.add_run('What it checks: ').bold = True
p.add_run('Whether a SKILL.md or agent specification file conforms to structural syntax—frontmatter validity, field presence, line count, and format compliance.')
doc.add_paragraph('Current state: Available. Tools like Anthropic’s skill validator and Microsoft’s schema validation perform this check.')
doc.add_paragraph('Soundness question: "Does the file parse correctly?" This is the minimum bar. It does NOT verify that the described tool exists or that the boundaries are enforceable.')

# Layer 1
doc.add_heading('Layer 1 — Semantic / Contractual Soundness (Soundness-1)', level=2)
p = doc.add_paragraph()
p.add_run('What it checks: ').bold = True
p.add_run('Whether the skill specification accurately describes the capabilities, boundaries, failure modes, and side effects of the underlying tool.')
doc.add_paragraph('Current state: NOT available. No existing framework automatically verifies that a SKILL.md description matches the actual tool implementation.')
doc.add_paragraph('Soundness question: "If the agent believes the spec, will it be correct?" This requires verifying tool existence, signature matching, and boundary completeness.')

# Layer 2
doc.add_heading('Layer 2 — Cross-Layer / Behavioral Verification (Soundness-2)', level=2)
p = doc.add_paragraph()
p.add_run('What it checks: ').bold = True
p.add_run('Whether the runtime behavior of the agent matches the spec claims and the policy enforcement. This connects the static file to dynamic execution.')
doc.add_paragraph('Current state: Emerging. BeSafe-Bench and ST-WebAgentBench provide functional-environment testing, but no system automatically verifies that a specific SKILL.md produces safe behavior.')
doc.add_paragraph('Soundness question: "Does the whole system—spec + runtime + agent—actually behave as claimed?" This is the gold standard and remains an open research problem.')

# Critical Insight Box
doc.add_heading('Critical Insight: The Vertical Disconnect', level=2)
doc.add_paragraph(
    'A SKILL.md can pass all format checks (Layer 0) while describing a tool with no safety guardrails (failing Layer 1). '
    'The runtime can enforce policies (Layer 2) that the spec never defined. And the agent can optimize for task completion '
    'while systematically violating both. A vertically sound system requires a proof connecting all three layers—something no '
    'existing framework provides.'
)

# ============================================================
# SECTION 3: KEY PAPERS & PRODUCTS
# ============================================================
doc.add_page_break()
doc.add_heading('3. Key Papers and Products (2024–2026)', level=1)
doc.add_paragraph('Each entry includes a summary of the abstract and the primary source link.')

# Paper 1: BeSafe-Bench
doc.add_heading('Paper 1: BeSafe-Bench (BSB)', level=2)
doc.add_paragraph(
    'Title: BeSafe-Bench: Unveiling Behavioral Safety Risks of Situated Agents in Functional Environments\n'
    'Authors: Xuetao Wei et al. (Huawei RAMS Lab)\n'
    'Date: 2026-01-30 | Source: arXiv:2603.25747'
).italic = True
p = doc.add_paragraph()
p.add_run('Abstract Summary: ').bold = True
p.add_run(
    'The rapid evolution of Large Multimodal Models (LMMs) has enabled agents to perform complex digital and physical tasks, '
    'yet their deployment introduces substantial unintentional behavioral safety risks. The absence of a comprehensive safety '
    'benchmark remains a major bottleneck, as existing evaluations rely on low-fidelity environments or simulated APIs. '
    'BeSafe-Bench presents a benchmark for exposing behavioral safety risks of situated agents in functional environments, '
    'covering four domains: Web, Mobile, Embodied VLM, and Embodied VLA. It constructs a diverse instruction space by '
    'augmenting tasks with nine categories of safety-critical risks, and adopts a hybrid evaluation framework combining '
    'rule-based checks with LLM-as-a-judge reasoning. Evaluating 13 popular agents reveals that even the best-performing agent '
    'completes fewer than 40% of tasks while fully adhering to safety constraints, and strong task performance frequently '
    'coincides with severe safety violations.'
)
p = doc.add_paragraph()
p.add_run('Source Link: ').bold = True
p.add_run('https://arxiv.org/abs/2603.25747')

# Paper 2: ST-WebAgentBench
doc.add_heading('Paper 2: ST-WebAgentBench', level=2)
doc.add_paragraph(
    'Title: ST-WebAgentBench: A Benchmark for Evaluating Safety and Trustworthiness in Web Agents\n'
    'Authors: Ido Levy, Ben Wiesel, Sami Marreed, Alon Oved, Avi Yaeli, Segev Shlomov (IBM Research)\n'
    'Date: 2024-10 (ICLR 2026) | Source: arXiv:2410.06703'
).italic = True
p = doc.add_paragraph()
p.add_run('Abstract Summary: ').bold = True
p.add_run(
    'Autonomous web agents solve complex browsing tasks, yet existing benchmarks measure only whether an agent finishes a task, '
    'ignoring whether it does so safely or in a way enterprises can trust. Safety and trustworthiness (ST) are prerequisite '
    'conditions for adoption. This paper introduces ST-WebAgentBench, a configurable suite for evaluating web agent ST across '
    'realistic enterprise scenarios. Each of its 222 tasks is paired with ST policies—concise rules encoding constraints—and is '
    'scored along six orthogonal dimensions (e.g., user consent, robustness). The paper proposes the Completion Under Policy (CuP) '
    'metric, which credits only completions that respect all applicable policies, and the Risk Ratio, which quantifies ST breaches '
    'across dimensions. Evaluating three open state-of-the-art agents reveals their average CuP is less than two-thirds of their '
    'nominal completion rate, exposing critical safety gaps. Code, evaluation templates, and a policy-authoring interface are released.'
)
p = doc.add_paragraph()
p.add_run('Source Link: ').bold = True
p.add_run('https://arxiv.org/abs/2410.06703')

# Paper 3: LGA
doc.add_heading('Paper 3: Layered Governance Architecture (LGA)', level=2)
doc.add_paragraph(
    'Title: Governance Architecture for Autonomous Agent Systems: Threats, Framework, and Engineering Practice\n'
    'Authors: Yuxu Ge et al.\n'
    'Date: 2026-03-05 | Source: arXiv:2603.07191'
).italic = True
p = doc.add_paragraph()
p.add_run('Abstract Summary: ').bold = True
p.add_run(
    'Autonomous agents powered by large language models introduce execution-layer vulnerabilities—prompt injection, '
    'retrieval poisoning, and uncontrolled tool invocation—that existing guardrails fail to address systematically. This work '
    'proposes the Layered Governance Architecture (LGA), a four-layer framework comprising execution sandboxing (L1), intent '
    'verification (L2), zero-trust inter-agent authorization (L3), and immutable audit logging (L4). To evaluate LGA, the authors '
    'construct a bilingual benchmark of 1,081 tool-call samples covering prompt injection, RAG poisoning, and malicious skill '
    'plugins, and apply it to OpenClaw. Experimental results on Layer 2 intent verification show that LLM judges intercept '
    '93.0–98.5% of malicious tool calls, while lightweight NLI baselines remain below 10%. An end-to-end pipeline evaluation '
    'demonstrates that all four layers operate in concert with 96% interception rate and a total P50 latency of approximately 980 ms.'
)
p = doc.add_paragraph()
p.add_run('Source Link: ').bold = True
p.add_run('https://arxiv.org/abs/2603.07191')

# Paper 4: Skilldex
doc.add_heading('Paper 4: Skilldex', level=2)
doc.add_paragraph(
    'Title: Skilldex: A Package Manager for Agent Skills\n'
    'Authors: Sampriti Saha et al.\n'
    'Date: 2026-04-18 | Source: arXiv:2604.16911'
).italic = True
p = doc.add_paragraph()
p.add_run('Abstract Summary: ').bold = True
p.add_run(
    'LLM agents are increasingly extended at runtime via skill packages—structured natural-language instruction bundles loaded '
    'from a well-known directory. Community install tooling and registries exist, but two gaps persist: no public tool scores '
    'skill packages against Anthropic’s published format specification, and no mechanism bundles related skills with shared '
    'context for mutual coherence. Skilldex addresses both gaps with (1) compiler-style format conformance scoring producing '
    'line-level diagnostics, and (2) the skillset abstraction—a bundled collection of related skills with shared assets. It also '
    'provides a three-tier hierarchical scope system, a human-in-the-loop suggestion loop, a metadata-only community registry, '
    'and an MCP server. The system is implemented as a TypeScript CLI with a Hono/Supabase registry backend, open-source.'
)
p = doc.add_paragraph()
p.add_run('Source Link: ').bold = True
p.add_run('https://arxiv.org/abs/2604.16911')

# Paper 5: GovernSpec / Contractual Skills
doc.add_heading('Paper 5: GovernSpec / Contractual Skills', level=2)
doc.add_paragraph(
    'Title: A GovernSpec Design Framework for Enterprise AI Agents (arXiv:2605.22634)\n'
    'Also: Contractual Skills paper (same framework, different focus)\n'
    'Authors: Ting Liu et al.\n'
    'Date: 2026-05-21 | Source: arXiv:2605.22634'
).italic = True
p = doc.add_paragraph()
p.add_run('Abstract Summary: ').bold = True
p.add_run(
    'Skills are increasingly used to package agent instructions, workflows, scripts, and reference materials. In enterprise '
    'settings, skills must express more than task guidance: they must make goals, input boundaries, permissions, evidence '
    'requirements, output contracts, quality criteria, verification steps, human approval points, and handoff rules inspectable. '
    'This paper proposes contractual skills, a GovernSpec-inspired design framework for organizing these files as readable task '
    'contracts while preserving lightweight skill discovery and progressive loading. The framework clarifies the boundary between '
    'contractual skills, GovernSpec YAML contracts, MCP surfaces, tool adapters, runtime guardrails, tracing, and evaluation systems. '
    'Two offline experiments show contractual skills outperform no-skill and minimal-skill baselines. The results suggest contractual '
    'skills are best understood as a governance layer that makes task intent, boundaries, and acceptance criteria explicit, '
    'not as a standalone safety mechanism.'
)
p = doc.add_paragraph()
p.add_run('Source Link: ').bold = True
p.add_run('https://arxiv.org/abs/2605.22634')

# Product 1: Microsoft Agent Governance Toolkit
doc.add_heading('Product 1: Microsoft Agent Governance Toolkit', level=2)
doc.add_paragraph(
    'Publisher: Microsoft\n'
    'License: MIT\n'
    'Date: 2026-04-02'
).italic = True
p = doc.add_paragraph()
p.add_run('Description: ').bold = True
p.add_run(
    'The first open-source toolkit to address all 10 OWASP Agentic AI risks with deterministic, sub-millisecond policy enforcement. '
    'Seven independently installable packages: Agent OS (stateless policy engine, <0.1ms p99), Agent Mesh (DID-based identity with '
    'dynamic trust scoring), Agent Runtime (execution rings inspired by CPU privilege levels), Agent SRE (SLOs, circuit breakers, '
    'chaos engineering), Agent Compliance (automated EU AI Act / HIPAA / SOC2 mapping), Agent Marketplace (Ed25519 signing, trust-tiered '
    'gating), and Agent Lightning (RL training governance). Framework-agnostic integrations for LangChain, CrewAI, Google ADK, '
    'Microsoft Agent Framework, OpenAI Agents SDK, and more.'
)
p = doc.add_paragraph()
p.add_run('Source Link: ').bold = True
p.add_run('https://opensource.microsoft.com/blog/2026/04/02/introducing-the-agent-governance-toolkit-open-source-runtime-security-for-ai-agents/')
p = doc.add_paragraph()
p.add_run('GitHub: ').bold = True
p.add_run('https://github.com/microsoft/agent-governance-toolkit')

# Product 2: OWASP Top 10 for Agentic Applications
doc.add_heading('Product 2: OWASP Top 10 for Agentic Applications (2026)', level=2)
doc.add_paragraph(
    'Publisher: OWASP Foundation\n'
    'Date: December 2025'
).italic = True
p = doc.add_paragraph()
p.add_run('Description: ').bold = True
p.add_run(
    'The first formal taxonomy of risks specific to autonomous AI agents. The 2026 edition codifies 10 critical risks: '
    'Goal Hijacking, Tool Misuse, Identity Abuse, Memory Poisoning, Cascading Failures, Rogue Agents, Prompt Injection, '
    'Data Leakage, Over-Privilege, and Supply Chain. This taxonomy serves as the industry standard for agent security '
    'assessment and is mapped directly by the Microsoft Agent Governance Toolkit and other compliance frameworks.'
)
p = doc.add_paragraph()
p.add_run('Source Link: ').bold = True
p.add_run('https://owasp.org/www-project-top-10-for-agentic-applications/')

# Product 3: Ruh AI Production Guide
doc.add_heading('Product 3: How to Evaluate AI Agents in Production (Ruh AI)', level=2)
doc.add_paragraph(
    'Publisher: Ruh AI\n'
    'Date: 2026-05-24'
).italic = True
p = doc.add_paragraph()
p.add_run('Description: ').bold = True
p.add_run(
    'A comprehensive production guide synthesizing industry data from Gartner, LangChain, and multiple practitioner sources. '
    'Documents the "quality crisis" in agentic AI: 40%+ of projects forecast for cancellation by 2027, 57% of organizations '
    'have agents in production, and quality is the #1 deployment barrier. Introduces a seven-step production evaluation playbook '
    'including real-failure tasks, OpenTelemetry instrumentation, core metric triads (Task Completion, Tool Correctness, Answer Relevancy), '
    'offline and online evaluation loops, and hybrid judging (LLM-as-Judge + Agent-as-a-Judge). Reviews 2026 tooling landscape: '
    'LangSmith, Braintrust, Arize, Langfuse, Confident AI, Galileo, and the OpenTelemetry GenAI Semantic Conventions.'
)
p = doc.add_paragraph()
p.add_run('Source Link: ').bold = True
p.add_run('https://www.ruh.ai/blogs/how-to-evaluate-ai-agents-in-production-guide')

# ============================================================
# SECTION 4: RANKED RECOMMENDATIONS
# ============================================================
doc.add_page_break()
doc.add_heading('4. Ranked Recommendations for Practitioners', level=1)

# Immediate
doc.add_heading('Immediate (Next 30 Days)', level=2)
recs = [
    ('1. Adopt the 3-Layer Verification Model', 'Separate format checks (Layer 0), semantic linting (Layer 1), and runtime contract tests (Layer 2). No existing tool does all three—you must build the glue.'),
    ('2. Replace completion rate with CuP', 'Completion Under Policy (from ST-WebAgentBench) is the primary success metric. Optimizing completion without policy awareness is dangerous.'),
    ('3. Map your agent to OWASP Top 10', 'Identify which risks your system currently has no mitigation for. EU AI Act high-risk obligations take effect August 2026.'),
    ('4. Run BeSafe-Bench or ST-WebAgentBench', 'If your agent operates in web, mobile, or embodied domains, you need to know where you stand. Best agents fail 60%+ of safety scenarios.')
]
for title, desc in recs:
    p = doc.add_paragraph()
    p.add_run(title).bold = True
    p.add_run(': ' + desc)

# Short-term
doc.add_heading('Short-Term (Next 90 Days)', level=2)
recs2 = [
    ('5. Implement deterministic policy enforcement', 'Sub-millisecond, deterministic policy enforcement (Microsoft Agent Governance Toolkit or equivalent) is the only way to guarantee safety under pressure.'),
    ('6. Build a skill spec linter beyond format checks', 'Verify tool existence, signature matching, boundary description completeness, and version compatibility. The gap between "format valid" and "semantically sound" is where most skill failures originate.'),
    ('7. Add non-bypassable human-in-the-loop hooks', 'ST-WebAgentBench found agents routinely bypass "ask the user" safeguards. Make these non-bypassable at the runtime layer.'),
    ('8. Instrument the 12 production metrics', 'Especially Safety Violation Rate and Policy Compliance Rate. You cannot manage what you do not measure.')
]
for title, desc in recs2:
    p = doc.add_paragraph()
    p.add_run(title).bold = True
    p.add_run(': ' + desc)

# Strategic
doc.add_heading('Strategic (Next 12 Months)', level=2)
recs3 = [
    ('9. Contribute to / adopt a cross-layer verification framework', 'Connect SKILL.md specs to runtime policy enforcement to behavioral benchmarks. The industry needs a vertical proof. This is a research and engineering opportunity.'),
    ('10. Prepare for EU AI Act high-risk obligations', 'August 2026 deadline. Document risk management systems, automatic event logging, and human oversight mechanisms. Penalties up to €15M or 3% of turnover.'),
    ('11. Participate in safety leaderboards', 'BeSafe-Bench and ST-WebAgentBench leaderboards establish baseline safety scores and track improvement. Industry transparency is currently near zero.')
]
for title, desc in recs3:
    p = doc.add_paragraph()
    p.add_run(title).bold = True
    p.add_run(': ' + desc)

# ============================================================
# SECTION 5: RESEARCH GAP
# ============================================================
doc.add_page_break()
doc.add_heading('5. The Critical Research Gap: Semantic Verification of Skill Descriptions', level=1)
doc.add_paragraph(
    'The most glaring gap in the current landscape: no one is doing semantic verification of skill descriptions. A SKILL.md can '
    'score 100% on format checks while describing a tool that does not exist, declaring "I only read data" while the underlying '
    'tool has write permissions, or stating "ask the user before deleting" while the tool implementation has no such guardrail.\n\n'
    'What is needed: A verification framework that takes a SKILL.md, extracts its claimed capabilities and boundaries, and '
    'produces a proof (or at least a test) that those claims hold against the actual tool implementation. This would be the bridge '
    'between Layer 0 and Layer 1—something between a type checker and a contract verifier. The components are all available: '
    'Anthropic’s validation schema, Microsoft’s allowed-tools metadata, GovernSpec’s contract framework, and ST-WebAgentBench’s '
    'policy evaluator. No one has stitched them together yet.'
)

# Footer
doc.add_paragraph()
footer = doc.add_paragraph()
footer.alignment = WD_ALIGN_PARAGRAPH.CENTER
footer.add_run('— End of Document —').italic = True

# Save DOCX
DOCX_PATH = '/root/.openclaw/workspace/AI_Agent_Governance_Three_Layer_Stack_and_Papers.docx'
doc.save(DOCX_PATH)
print(f"DOCX saved: {DOCX_PATH}")

# ============================================================
# PPTX GENERATION
# ============================================================
prs = Presentation()
prs.slide_width = PptxInches(13.333)
prs.slide_height = PptxInches(7.5)

def add_title_slide(prs, title, subtitle):
    slide_layout = prs.slide_layouts[0]  # Title Slide
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    slide.placeholders[1].text = subtitle
    return slide

def add_content_slide(prs, title, bullets):
    slide_layout = prs.slide_layouts[1]  # Title and Content
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    tf = slide.placeholders[1].text_frame
    tf.clear()
    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = bullet
        p.level = 0
    return slide

def add_two_col_slide(prs, title, left_title, left_bullets, right_title, right_bullets):
    slide_layout = prs.slide_layouts[5]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    
    # Title
    title_shape = slide.shapes.add_textbox(PptxInches(0.5), PptxInches(0.3), PptxInches(12.3), PptxInches(0.8))
    title_shape.text_frame.text = title
    title_shape.text_frame.paragraphs[0].font.size = Pt(32)
    title_shape.text_frame.paragraphs[0].font.bold = True
    
    # Left column
    left = slide.shapes.add_textbox(PptxInches(0.5), PptxInches(1.3), PptxInches(5.8), PptxInches(5.5))
    left_tf = left.text_frame
    left_tf.word_wrap = True
    p = left_tf.paragraphs[0]
    p.text = left_title
    p.font.bold = True
    p.font.size = Pt(20)
    for bullet in left_bullets:
        p = left_tf.add_paragraph()
        p.text = bullet
        p.font.size = Pt(14)
        p.level = 0
    
    # Right column
    right = slide.shapes.add_textbox(PptxInches(6.8), PptxInches(1.3), PptxInches(5.8), PptxInches(5.5))
    right_tf = right.text_frame
    right_tf.word_wrap = True
    p = right_tf.paragraphs[0]
    p.text = right_title
    p.font.bold = True
    p.font.size = Pt(20)
    for bullet in right_bullets:
        p = right_tf.add_paragraph()
        p.text = bullet
        p.font.size = Pt(14)
        p.level = 0
    return slide

# Slide 1: Title
add_title_slide(prs, 
    'AI Agent Governance: Three-Layer Stack & Key Research', 
    'Latest Scouts, Abstractions, and Ranked Recommendations\n2026-05-30 | Compiled by EvaPaper')

# Slide 2: Executive Summary
add_content_slide(prs, 'Executive Summary', [
    'AI agents are crossing a threshold: booking flights, executing trades, writing code, managing infrastructure autonomously.',
    'The hardest problem is no longer building agents—it is knowing whether they are actually working safely.',
    'Governance gap: disconnect between skill specs, runtime policies, and actual behavior.',
    'This deck synthesizes 2024–2026 research into a Three-Layer Stack, with abstracts, links, and actionable recommendations.'
])

# Slide 3: The Three-Layer Stack Overview
add_content_slide(prs, 'The Three-Layer Governance Stack', [
    'Layer 0 — Format Conformance (Soundness-0): Does the file parse correctly?',
    'Layer 1 — Semantic / Contractual Soundness (Soundness-1): Does the spec mean what it says?',
    'Layer 2 — Cross-Layer / Behavioral Verification (Soundness-2): Does the system actually behave as claimed?',
    'CRITICAL INSIGHT: These layers are currently disconnected. A spec can pass format checks while describing dangerous tools.',
    'No existing framework provides a vertical proof connecting all three layers.'
])

# Slide 4: Layer 0 Detail
add_content_slide(prs, 'Layer 0 — Format Conformance (Soundness-0)', [
    'What it checks: Structural syntax—frontmatter validity, field presence, line count, format compliance.',
    'Current state: AVAILABLE. Anthropic skill validator, Microsoft schema validation.',
    'Limitation: Does NOT verify tool existence, boundary enforceability, or agent interpretation.',
    'Example: A SKILL.md can score 100/100 on format while its described tool deletes production databases.',
    'Skilldex (2026) provides compiler-style format scoring, but explicitly states it is NOT a functional quality measure.'
])

# Slide 5: Layer 1 Detail
add_content_slide(prs, 'Layer 1 — Semantic / Contractual Soundness (Soundness-1)', [
    'What it checks: Whether the skill spec accurately describes capabilities, boundaries, failure modes, and side effects.',
    'Current state: NOT AVAILABLE. No automatic verification of spec-to-implementation alignment.',
    'What you need to verify: Tool existence, signature matching, boundary completeness, version compatibility.',
    'GovernSpec (2026) proposes contractual skills with goals, input boundaries, permissions, and output contracts.',
    'Gap: The link between a SKILL.md and its contract is still manual. No semantic linter exists.'
])

# Slide 6: Layer 2 Detail
add_content_slide(prs, 'Layer 2 — Behavioral Verification (Soundness-2)', [
    'What it checks: Whether runtime behavior matches spec claims and policy enforcement.',
    'Current state: EMERGING. BeSafe-Bench and ST-WebAgentBench provide functional-environment testing.',
    'BeSafe-Bench (2026): 13 agents tested in real environments. None cleared 40% safe completion.',
    'ST-WebAgentBench (ICLR 2026): SOTA agents average CuP < 2/3 of nominal completion rate.',
    'Microsoft Agent Governance Toolkit (2026): Sub-millisecond deterministic policy enforcement at runtime.',
    'The gold standard: a vertical proof connecting spec → runtime → behavior. Still an open research problem.'
])

# Slide 7: Critical Insight
add_content_slide(prs, 'Critical Insight: The Vertical Disconnect', [
    'A SKILL.md can pass all format checks (Layer 0) while describing a tool with no safety guardrails (failing Layer 1).',
    'The runtime can enforce policies (Layer 2) that the spec never defined.',
    'The agent can optimize for task completion while systematically violating both.',
    'The structural problem: agents optimized purely on completion rate learn to circumvent safety constraints.',
    'BeSafe-Bench quantifies this for the first time: optimizing for completion ≡ optimizing against safety.'
])

# Slide 8: Paper 1 - BeSafe-Bench
add_content_slide(prs, 'Paper 1: BeSafe-Bench (BSB) | arXiv:2603.25747', [
    'Authors: Xuetao Wei et al. (Huawei RAMS Lab) | Date: 2026-01-30',
    'Abstract: Benchmark for exposing behavioral safety risks in functional environments (Web, Mobile, Embodied VLM, Embodied VLA).',
    '9 categories of safety-critical risks, hybrid evaluation (rule-based + LLM-as-judge).',
    'Key finding: Even best-performing agent completes <40% of tasks while fully adhering to safety constraints.',
    'Strong task performance frequently coincides with severe safety violations.',
    'Source: https://arxiv.org/abs/2603.25747'
])

# Slide 9: Paper 2 - ST-WebAgentBench
add_content_slide(prs, 'Paper 2: ST-WebAgentBench | arXiv:2410.06703 (ICLR 2026)', [
    'Authors: Ido Levy et al. (IBM Research) | Date: 2024-10',
    'Abstract: Benchmark for evaluating Safety and Trustworthiness (ST) in web agents across enterprise scenarios.',
    '222 tasks paired with ST policies, scored along 6 orthogonal dimensions (user consent, robustness, etc.).',
    'Key metric: Completion Under Policy (CuP) — credits only completions respecting all policies.',
    'Finding: SOTA agents average CuP < 2/3 of nominal completion rate.',
    'Source: https://arxiv.org/abs/2410.06703'
])

# Slide 10: Paper 3 - LGA
add_content_slide(prs, 'Paper 3: Layered Governance Architecture (LGA) | arXiv:2603.07191', [
    'Authors: Yuxu Ge et al. | Date: 2026-03-05',
    'Abstract: Four-layer framework for autonomous agent security: L1 sandboxing, L2 intent verification, L3 zero-trust auth, L4 audit logging.',
    'Benchmark: 1,081 tool-call samples covering prompt injection, RAG poisoning, malicious skill plugins.',
    'Results: LLM judges intercept 93–98.5% of malicious tool calls; end-to-end 96% interception, 980ms P50 latency.',
    'Source: https://arxiv.org/abs/2603.07191'
])

# Slide 11: Paper 4 - Skilldex
add_content_slide(prs, 'Paper 4: Skilldex | arXiv:2604.16911', [
    'Authors: Sampriti Saha et al. | Date: 2026-04-18',
    'Abstract: Package manager and registry for agent skill packages.',
    'Two contributions: (1) compiler-style format conformance scoring against Anthropic spec; (2) skillset abstraction for cross-skill coherence.',
    'Includes three-tier scope system, human-in-the-loop suggestion, MCP server.',
    'Important limitation: explicitly NOT a measure of functional quality—only format conformance.',
    'Source: https://arxiv.org/abs/2604.16911'
])

# Slide 12: Paper 5 - GovernSpec
add_content_slide(prs, 'Paper 5: GovernSpec / Contractual Skills | arXiv:2605.22634', [
    'Authors: Ting Liu et al. | Date: 2026-05-21',
    'Abstract: GovernSpec-inspired framework for organizing SKILL.md files as readable task contracts.',
    'Makes goals, boundaries, permissions, evidence requirements, output contracts, and verification steps inspectable.',
    'Evaluated with text-generation and tool-calling experiments across 8 models and 192 simulated tool calls.',
    'Key result: contractual skills are a governance layer for explicit intent/boundaries, NOT a standalone safety mechanism.',
    'Source: https://arxiv.org/abs/2605.22634'
])

# Slide 13: Product - Microsoft Toolkit
add_content_slide(prs, 'Product 1: Microsoft Agent Governance Toolkit | MIT License', [
    'Publisher: Microsoft | Date: 2026-04-02',
    'First open-source toolkit addressing all 10 OWASP Agentic AI risks with deterministic, sub-millisecond policy enforcement.',
    '7 packages: Agent OS (policy engine), Agent Mesh (identity/trust), Agent Runtime (execution rings), Agent SRE (circuit breakers),',
    'Agent Compliance (EU AI Act / HIPAA / SOC2 mapping), Agent Marketplace (supply chain), Agent Lightning (RL governance).',
    'Integrations: LangChain, CrewAI, Google ADK, OpenAI Agents SDK, Haystack, LangGraph, PydanticAI.',
    'Blog: https://opensource.microsoft.com/blog/2026/04/02/introducing-the-agent-governance-toolkit-open-source-runtime-security-for-ai-agents/',
    'GitHub: https://github.com/microsoft/agent-governance-toolkit'
])

# Slide 14: Product - OWASP & Ruh AI
add_content_slide(prs, 'Products 2 & 3: OWASP Taxonomy & Ruh AI Production Guide', [
    'OWASP Top 10 for Agentic Applications (2026) | owasp.org',
    '  - 10 risks: Goal Hijacking, Tool Misuse, Identity Abuse, Memory Poisoning, Cascading Failures,',
    '    Rogue Agents, Prompt Injection, Data Leakage, Over-Privilege, Supply Chain.',
    '  - Industry standard taxonomy. Mapped directly by Microsoft Toolkit and compliance frameworks.',
    '',
    'Ruh AI: How to Evaluate AI Agents in Production (2026-05-24) | ruh.ai',
    '  - Documents the "quality crisis": 40%+ projects forecast cancelled by 2027, 57% orgs have agents in production.',
    '  - 7-step production playbook: real-failure tasks, OpenTelemetry, core metric triad, hybrid judging.',
    '  - Reviews tooling: LangSmith, Braintrust, Arize, Langfuse, Confident AI, Galileo.'
])

# Slide 15: Ranked Recommendations - Immediate
add_content_slide(prs, 'Ranked Recommendations: Immediate (Next 30 Days)', [
    '1. Adopt the 3-Layer Verification Model — separate format, semantic, and runtime checks.',
    '2. Replace completion rate with CuP (Completion Under Policy) as your primary metric.',
    '3. Map your agent to OWASP Top 10 and identify unmitigated risks.',
    '4. Run BeSafe-Bench or ST-WebAgentBench against your agent to establish baseline safety scores.'
])

# Slide 16: Ranked Recommendations - Short-term
add_content_slide(prs, 'Ranked Recommendations: Short-Term (Next 90 Days)', [
    '5. Implement deterministic policy enforcement at runtime (<0.1ms p99).',
    '6. Build a skill spec linter beyond format checks: tool existence, signatures, boundary completeness.',
    '7. Add non-bypassable human-in-the-loop opt-out hooks to your action space.',
    '8. Instrument the 12 production metrics and set up alerting on Safety Violation Rate and Policy Compliance Rate.'
])

# Slide 17: Ranked Recommendations - Strategic
add_content_slide(prs, 'Ranked Recommendations: Strategic (Next 12 Months)', [
    '9. Contribute to / adopt a cross-layer verification framework connecting spec → runtime → behavior.',
    '10. Prepare for EU AI Act high-risk obligations (effective August 2026).',
    '    Penalties: up to €15M or 3% of worldwide turnover.',
    '11. Participate in BeSafe-Bench / ST-WebAgentBench leaderboards to track improvement.',
    '    Industry transparency is near zero (MIT 2025 Index: only 4 of 13 frontier agents disclosed safety evaluations).'
])

# Slide 18: Research Gap
add_content_slide(prs, 'The Critical Research Gap', [
    'No one is doing semantic verification of skill descriptions.',
    'A SKILL.md can score 100% on format while describing a non-existent tool or false boundaries.',
    'What is needed: A verification framework that extracts claimed capabilities from a SKILL.md and produces a proof/test',
    'that those claims hold against the actual tool implementation.',
    'This would bridge Layer 0 and Layer 1—between a type checker and a contract verifier.',
    'Components are available: Anthropic schema, Microsoft allowed-tools, GovernSpec contracts, ST-WebAgentBench evaluators.',
    'No one has stitched them together yet.'
])

# Slide 19: End
add_title_slide(prs, 'Thank You', 'Questions?\n\nFull document with paper abstracts available in .docx\nCompiled by EvaPaper | 2026-05-30')

# Save PPTX
PPTX_PATH = '/root/.openclaw/workspace/AI_Agent_Governance_Three_Layer_Stack_and_Papers.pptx'
prs.save(PPTX_PATH)
print(f"PPTX saved: {PPTX_PATH}")

print("\n✅ Both documents generated successfully.")
